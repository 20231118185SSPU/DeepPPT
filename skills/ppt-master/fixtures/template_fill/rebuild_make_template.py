#!/usr/bin/env python3
"""Regenerate the synthetic 4-slide template-fill source PPTX (content-equivalent).

Run from the repo root:
    python3 skills/ppt-master/fixtures/template_fill/rebuild_make_template.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "sources" / "acme_template.pptx"


def add_title_body(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tb.text_frame.text = title
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    body = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5))
    for i, b in enumerate(bullets):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
    return s


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_body(prs, "Q1 业务回顾（合成模板页）",
                   ["营收：¥1,234 万（同比 +12.4%）", "客户数：18,500（环比 +8.1%）", "核心产品：ACME Platform 3.0"])
    add_title_body(prs, "Q2 增长计划", ["华东区域拓展（3 城）", "渠道伙伴计划 v2", "研发投入占比提升至 18%"])
    add_title_body(prs, "风险与应对", ["供应链：双源策略", "合规：季度审计", "汇率：套期保值"])
    add_title_body(prs, "附录：口径说明", ["财务口径为合并报表", "同比基期为 2025 同期"])
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
