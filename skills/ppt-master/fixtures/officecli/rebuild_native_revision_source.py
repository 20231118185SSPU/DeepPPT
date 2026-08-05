#!/usr/bin/env python3
"""Regenerate the native-revision synthetic PPTX fixture.

Produces a 4-slide deck with text, table, chart, picture, and grouped shapes,
all with predictable stable IDs for atomic plan testing.

Run from the repo root:
    python skills/ppt-master/fixtures/officecli/rebuild_native_revision_source.py
"""

import io
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent / "native_revision_source.pptx"
REPO_ROOT = Path(__file__).resolve().parent.parent.parent.parent.parent


class _LenientZip(zipfile.ZipFile):
    """python-pptx chart embeds may carry a bad CRC; PowerPoint tolerates it."""

    def _update_crc(self, newdata, eof=False):  # type: ignore[override]
        pass


_FIXED_STAMP = "1980-01-01T00:00:00Z"


def _normalize_embedded_xlsx(data: bytes) -> bytes:
    """Neutralize python-pptx chart embed core.xml timestamps for determinism."""
    import io as _io
    import re

    with _LenientZip(_io.BytesIO(data)) as z:
        entries = {i.filename: z.read(i.filename) for i in z.infolist()}
    core = entries.get("docProps/core.xml", b"")
    if core:
        text = core.decode("utf-8")
        text = re.sub(
            r"(?<=\>)\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}Z(?=</dcterms:(?:created|modified)>)",
            _FIXED_STAMP,
            text,
        )
        entries["docProps/core.xml"] = text.encode("utf-8")
    buf = _io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as out:
        for name in sorted(entries):
            out.writestr(zipfile.ZipInfo(name, (1980, 1, 1, 0, 0, 0)), entries[name])
    return buf.getvalue()


def _deterministic_rewrite(path: Path) -> None:
    """Rezip with fixed timestamps/attrs so repeated rebuilds are byte-identical."""
    tmp = path.with_name(path.name + ".tmp")
    with _LenientZip(path) as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx":
                data = _normalize_embedded_xlsx(data)
            item.date_time = (1980, 1, 1, 0, 0, 0)
            item.create_system = 0
            item.external_attr = 0
            item.compress_type = zipfile.ZIP_DEFLATED
            item.extra = b""
            item.comment = b""
            out.writestr(item, data)
    path.unlink()
    tmp.replace(path)


def _add_title(slide, text, top=Inches(0.3), left=Inches(0.5)):
    """Add a consistent title shape."""
    txBox = slide.shapes.add_textbox(left, top, Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    return txBox


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- Slide 1: Text shapes (title + 3 text boxes) ----
    s1 = prs.slides.add_slide(blank)
    _add_title(s1, "Page 1 - Text Shapes")
    tb1 = s1.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.5), Inches(2.5))
    tb1.text_frame.text = "Text Block A\n\nThis is the first text block.\nIt contains multiple lines\nof synthetic content."
    tb1.text_frame.paragraphs[0].font.size = Pt(18)

    tb2 = s1.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(2.5))
    tb2.text_frame.text = "Text Block B\n\nThis is the second text block.\nEach block has a stable ID\nfor targeted mutation tests."
    tb2.text_frame.paragraphs[0].font.size = Pt(18)

    tb3 = s1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(1.0))
    tb3.text_frame.text = "Footer: Synthetic fixture — no sensitive content."
    tb3.text_frame.paragraphs[0].font.size = Pt(12)
    tb3.text_frame.paragraphs[0].font.italic = True

    # ---- Slide 2: Table ----
    s2 = prs.slides.add_slide(blank)
    _add_title(s2, "Page 2 - Data Table")
    table_shape = s2.shapes.add_table(5, 4, Inches(0.5), Inches(1.6), Inches(12), Inches(5))
    table = table_shape.table
    headers = ["Product", "Q1", "Q2", "Q3"]
    rows = [
        ["Widget A", "120", "135", "148"],
        ["Widget B", "89", "92", "97"],
        ["Gadget X", "45", "52", "61"],
        ["Total", "254", "279", "306"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        for p in table.cell(0, c).text_frame.paragraphs:
            p.font.bold = True
    for r, row in enumerate(rows):
        for c, v in enumerate(row):
            table.cell(r + 1, c).text = v

    # ---- Slide 3: Chart + Picture ----
    s3 = prs.slides.add_slide(blank)
    _add_title(s3, "Page 3 - Chart and Picture")

    # Chart (bar chart)
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Widget A", (120, 135, 148))
    chart_data.add_series("Widget B", (89, 92, 97))
    chart_shape = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.7), Inches(6), Inches(4.2),
        chart_data,
    )

    # Picture (synthetic)
    img = Image.new("RGB", (400, 225), (52, 120, 180))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    pic_shape = s3.shapes.add_picture(buf, Inches(7.5), Inches(1.7), Inches(5), Inches(2.8))
    pic_label = s3.shapes.add_textbox(Inches(7.5), Inches(4.8), Inches(5), Inches(0.4))
    pic_label.text_frame.text = "Fig.1 — Synthetic test image"
    pic_label.text_frame.paragraphs[0].font.size = Pt(11)

    # ---- Slide 4: Grouped shapes (temporary, for group testing) ----
    s4 = prs.slides.add_slide(blank)
    _add_title(s4, "Page 4 - Grouped Shapes")

    # Three rectangles that form a "card" group
    rect1 = s4.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.5), Inches(1.6), Inches(3.8), Inches(2.5),
    )
    rect1.text_frame.text = "Card A"
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = RGBColor(200, 220, 240)

    rect2 = s4.shapes.add_shape(
        1,
        Inches(4.8), Inches(1.6), Inches(3.8), Inches(2.5),
    )
    rect2.text_frame.text = "Card B"
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = RGBColor(220, 200, 240)

    rect3 = s4.shapes.add_shape(
        1,
        Inches(9.1), Inches(1.6), Inches(3.8), Inches(2.5),
    )
    rect3.text_frame.text = "Card C"
    rect3.fill.solid()
    rect3.fill.fore_color.rgb = RGBColor(240, 220, 200)

    # Text below cards
    s4_tb = s4.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12), Inches(1.5))
    s4_tb.text_frame.text = (
        "Three colored cards with text labels.\n"
        "Each card has a stable shape ID.\n"
        "Group testing verifies unaddressed objects\n"
        "remain unchanged after atomic batch."
    )
    for p in s4_tb.text_frame.paragraphs:
        p.font.size = Pt(14)

    # Save
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    _deterministic_rewrite(OUT)
    print(f"Written {OUT}")
    print(f"  Slides: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
