#!/usr/bin/env python3
"""Regenerate the synthetic 3-slide native-enhance source PPTX (content-equivalent).

Run from the repo root:
    python3 skills/ppt-master/fixtures/enhance/rebuild_make_source.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "sources" / "enhance_source.pptx"


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for title, body in [
        ("Enhance 合成页 1", ["要点 A", "要点 B", "要点 C"]),
        ("Enhance 合成页 2", ["数据 1", "数据 2"]),
        ("Enhance 合成页 3", ["结语与致谢"]),
    ]:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
        tb.text_frame.text = title
        b = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5))
        for i, t in enumerate(body):
            p = b.text_frame.paragraphs[0] if i == 0 else b.text_frame.add_paragraph()
            p.text = t
            p.font.size = Pt(18)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
