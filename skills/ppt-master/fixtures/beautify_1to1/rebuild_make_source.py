#!/usr/bin/env python3
"""Regenerate the synthetic 3-slide beautify source PPTX (content-equivalent).

Run from the repo root:
    python3 skills/ppt-master/fixtures/beautify_1to1/rebuild_make_source.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "sources" / "beautify_source.pptx"


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for title, body in [
        ("第一章 背景与动机", ["行业现状", "问题定义", "目标"]),
        ("第二章 方案设计", ["架构总览", "核心模块", "接口约定"]),
        ("第三章 验证与结论", ["实验设置", "结果对比", "总结"]),
    ]:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].font.size = Pt(28)
        b = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5))
        for i, t in enumerate(body):
            p = b.text_frame.paragraphs[0] if i == 0 else b.text_frame.add_paragraph()
            p.text = t
            p.font.size = Pt(18)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
