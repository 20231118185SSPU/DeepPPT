#!/usr/bin/env python3
"""Regenerate the complex PPTX fidelity fixture (merged table/image/notes/symbols).

Run from the repo root:
    python3 skills/ppt-master/fixtures/pptx_complex/rebuild_make_source.py
"""
import io
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

OUT = Path(__file__).resolve().parent / "complex_source.pptx"


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 页 1：合并单元格表格
    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tb.text_frame.text = "合成表格页"
    table = s1.shapes.add_table(4, 3, Inches(0.5), Inches(1.6), Inches(8), Inches(3)).table
    data = [["指标", "2025", "2026"], ["营收(万元)", "1,234", "1,386"],
            ["增速", "12.4%", "12.3%"], ["客户数", "18,500", "19,912"]]
    for r, row in enumerate(data):
        for c, v in enumerate(row):
            table.cell(r, c).text = v
    table.cell(1, 0).merge(table.cell(3, 0))
    table.cell(1, 0).text = "核心指标"

    # 页 2：图片 + 备注
    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tb2.text_frame.text = "图片与备注页"
    img = Image.new("RGB", (320, 180), (22, 50, 79))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    s2.shapes.add_picture(buf, Inches(0.5), Inches(1.6), Inches(4), Inches(2.25))
    s2.notes_slide.notes_text_frame.text = "本页备注：图片为合成占位图，无敏感内容。"

    # 页 3：符号文本
    s3 = prs.slides.add_slide(blank)
    tb3 = s3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.9))
    tb3.text_frame.text = "符号页"
    b3 = s3.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5))
    b3.text_frame.text = "圈号 ①②③ · 箭头 → · 检查 ✓ · 温度 ±2℃ · 汇率 ¥7.12"

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
