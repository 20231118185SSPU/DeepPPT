#!/usr/bin/env python3
"""
PPT Master - End-to-End Pipeline Validation

Validates that a completed project's deliverables are internally consistent:
SVG count matches spec_lock page count, every slide has speaker notes,
all declared images exist, and the exported PPTX is structurally sound.

Usage:
    python3 scripts/e2e_validate.py <project_path>
    python3 scripts/e2e_validate.py <project_path> --pptx exports/deck.pptx
    python3 scripts/e2e_validate.py <project_path> --strict

Examples:
    python3 scripts/e2e_validate.py projects/my_deck_ppt169_20260626
    python3 scripts/e2e_validate.py projects/my_deck_ppt169_20260626 --pptx exports/my_deck.pptx

Dependencies:
    python-pptx (only when --pptx is provided; skipped otherwise)

Notes:
    - Checks run without --pptx: page count, SVG files, speaker notes, image files
    - Additional checks with --pptx: PPTX opens, slide count matches, notes embedded
    - Exit 0 = all checks pass; exit 1 = errors found; exit 2 = usage error
"""

import sys
import argparse
import re
from pathlib import Path
from typing import Optional


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def parse_page_ids_from_spec_lock(spec_lock_path: Path) -> list[str]:
    """Extract page IDs (P01, P02, ...) from spec_lock.md's page_rhythm section.

    Supports both table format (| P01 | anchor |) and list format (- P01: anchor).
    """
    text = spec_lock_path.read_text(encoding="utf-8")
    ids: list[str] = []
    in_rhythm = False
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.lower().startswith("## page_rhythm"):
            in_rhythm = True
            continue
        if in_rhythm and stripped.startswith("## "):
            break
        if not in_rhythm:
            continue
        # Table format: | P01 | anchor | ...
        m = re.match(r"^\|\s*(P\d+)\s*\|", stripped)
        if m:
            ids.append(m.group(1))
            continue
        # List format: - P01: anchor
        m = re.match(r"^[-*]\s*(P\d+)\s*:", stripped)
        if m:
            ids.append(m.group(1))
    return ids


def parse_image_refs_from_spec_lock(spec_lock_path: Path) -> list[str]:
    """Extract image filenames from spec_lock.md's images section.

    Supports formats:
      - images/image_001.png
      - images/image_001.png | no-crop
      - P01_cover_bg: images/P01_cover_bg.png
      - label: images/foo.png | no-crop
    """
    text = spec_lock_path.read_text(encoding="utf-8")
    refs: list[str] = []
    in_images = False
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.lower().startswith("## images"):
            in_images = True
            continue
        if in_images and stripped.startswith("## "):
            break
        if not in_images or not stripped.startswith("- "):
            continue
        # Strip leading "- "
        content = stripped[2:].split("|")[0].split("#")[0].strip()
        if not content:
            continue
        # Handle "label: images/foo.png" → extract the path after ":"
        if ":" in content:
            path_part = content.split(":", 1)[1].strip()
        else:
            path_part = content
        if path_part and ("images/" in path_part or path_part.endswith(".png")
                          or path_part.endswith(".jpg") or path_part.endswith(".jpeg")):
            refs.append(path_part)
    return refs


def find_note_file(notes_dir: Path, page_id: str) -> Optional[Path]:
    """Find a speaker notes .md file matching the page ID."""
    if not notes_dir.is_dir():
        return None
    # Notes files are named like P01_cover.md, P02_toc.md, etc.
    for f in notes_dir.iterdir():
        if f.suffix == ".md" and f.stem != "total" and f.name.startswith(page_id):
            return f
    return None


def count_pptx_slides(pptx_path: Path) -> tuple[int, list[str]]:
    """Open a PPTX and return (slide_count, list_of_notes_per_slide).

    Returns notes as a list of strings (empty string if no notes on a slide).
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError(
            "python-pptx is not installed. Install it with: pip install python-pptx"
        )

    prs = Presentation(str(pptx_path))
    notes_list: list[str] = []
    for slide in prs.slides:
        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
        notes_list.append(notes_text)
    return len(prs.slides), notes_list


# ---------------------------------------------------------------------------
# Checks
# ---------------------------------------------------------------------------

def _safe_print(msg: str) -> None:
    """Print with encoding fallback for Windows consoles."""
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode("utf-8", errors="replace").decode("utf-8", errors="replace"))


class CheckResult:
    """Accumulates pass/fail/warn messages."""

    def __init__(self) -> None:
        self.passes: list[str] = []
        self.warns: list[str] = []
        self.errors: list[str] = []

    def ok(self, msg: str) -> None:
        self.passes.append(msg)

    def warn(self, msg: str) -> None:
        self.warns.append(msg)

    def fail(self, msg: str) -> None:
        self.errors.append(msg)

    @property
    def has_errors(self) -> bool:
        return len(self.errors) > 0

    def report(self) -> None:
        for m in self.passes:
            _safe_print(f"  [PASS] {m}")
        for m in self.warns:
            _safe_print(f"  [WARN] {m}")
        for m in self.errors:
            _safe_print(f"  [FAIL] {m}")


def check_page_count(project: Path, result: CheckResult) -> list[str]:
    """Check spec_lock page count vs SVG file count."""
    spec_lock = project / "spec_lock.md"
    if not spec_lock.is_file():
        result.fail("spec_lock.md not found — cannot determine expected page count")
        return []

    page_ids = parse_page_ids_from_spec_lock(spec_lock)
    if not page_ids:
        result.fail("No pages found in spec_lock.md page_rhythm section")
        return []

    svg_dir = project / "svg_output"
    svg_files = sorted(svg_dir.glob("P*.svg")) if svg_dir.is_dir() else []
    svg_stems = [f.stem for f in svg_files]

    if len(svg_files) == len(page_ids):
        result.ok(f"SVG count ({len(svg_files)}) = spec_lock page count ({len(page_ids)})")
    else:
        result.fail(
            f"SVG count ({len(svg_files)}) ≠ spec_lock page count ({len(page_ids)}). "
            f"Missing: {set(page_ids) - set(svg_stems)}, "
            f"Extra: {set(svg_stems) - set(page_ids)}"
        )

    return svg_stems


def check_speaker_notes(project: Path, svg_stems: list[str], result: CheckResult) -> None:
    """Check that every SVG page has a corresponding speaker notes file."""
    notes_dir = project / "notes"
    total_md = notes_dir / "total.md"

    if not notes_dir.is_dir():
        result.fail("notes/ directory not found")
        return

    if total_md.is_file():
        result.ok("notes/total.md exists (master speaker notes)")
    else:
        result.warn("notes/total.md not found — may not have been generated yet")

    missing: list[str] = []
    for stem in svg_stems:
        page_id = stem.split("_")[0]  # "P01_cover" → "P01"
        note_file = find_note_file(notes_dir, page_id)
        if note_file is None:
            missing.append(stem)

    if not missing:
        result.ok(f"All {len(svg_stems)} pages have speaker notes files")
    else:
        result.warn(
            f"{len(missing)}/{len(svg_stems)} pages missing speaker notes: "
            + ", ".join(missing)
        )


def check_images(project: Path, result: CheckResult) -> None:
    """Check that all images declared in spec_lock exist on disk."""
    spec_lock = project / "spec_lock.md"
    if not spec_lock.is_file():
        return

    refs = parse_image_refs_from_spec_lock(spec_lock)
    if not refs:
        result.ok("No images declared in spec_lock (text-only deck)")
        return

    images_dir = project / "images"
    missing: list[str] = []
    for ref in refs:
        # ref is like "images/image_001.png" — resolve relative to project
        full_path = project / ref
        if not full_path.is_file():
            missing.append(ref)

    if not missing:
        result.ok(f"All {len(refs)} declared images exist on disk")
    else:
        result.fail(
            f"{len(missing)}/{len(refs)} declared images missing: "
            + ", ".join(missing[:5])
            + (" ..." if len(missing) > 5 else "")
        )


def check_pptx(pptx_path: Path, expected_pages: int, result: CheckResult) -> None:
    """Validate the exported PPTX file."""
    if not pptx_path.is_file():
        result.fail(f"PPTX not found: {pptx_path}")
        return

    try:
        slide_count, notes_list = count_pptx_slides(pptx_path)
    except RuntimeError as e:
        result.fail(str(e))
        return
    except Exception as e:
        result.fail(f"PPTX cannot be opened: {e}")
        return

    result.ok(f"PPTX opens successfully ({pptx_path.name})")

    if expected_pages > 0 and slide_count == expected_pages:
        result.ok(f"PPTX slide count ({slide_count}) = expected ({expected_pages})")
    elif expected_pages > 0:
        result.fail(f"PPTX slide count ({slide_count}) ≠ expected ({expected_pages})")
    else:
        result.ok(f"PPTX contains {slide_count} slides")

    with_notes = sum(1 for n in notes_list if n)
    without_notes = slide_count - with_notes
    if without_notes == 0:
        result.ok(f"All {slide_count} slides have speaker notes embedded")
    elif without_notes <= 2:
        result.warn(
            f"{without_notes}/{slide_count} slides missing embedded notes "
            f"(indices: {[i for i, n in enumerate(notes_list) if not n]})"
        )
    else:
        result.warn(
            f"{without_notes}/{slide_count} slides missing embedded notes"
        )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="End-to-end pipeline validation for PPT Master projects.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("project_path", help="Path to the project directory")
    parser.add_argument(
        "--pptx",
        help="Path to the exported PPTX file (relative to project or absolute). "
        "When provided, validates PPTX structure and embedded notes.",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Treat warnings as errors (exit 2 if any warnings).",
    )
    return parser


def main(argv: Optional[list[str]] = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    project = Path(args.project_path)
    if not project.is_dir():
        print(f"Error: project directory not found: {project}", file=sys.stderr)
        return 2

    print(f"Validating: {project.name}\n")

    result = CheckResult()

    # 1. Page count: spec_lock vs SVGs
    print("Page Count")
    svg_stems = check_page_count(project, result)
    print()

    # 2. Speaker notes
    print("Speaker Notes")
    check_speaker_notes(project, svg_stems, result)
    print()

    # 3. Image completeness
    print("Images")
    check_images(project, result)
    print()

    # 4. PPTX validation (optional)
    if args.pptx:
        pptx_path = Path(args.pptx)
        if not pptx_path.is_absolute():
            pptx_path = project / pptx_path

        print("PPTX Export")
        expected = len(svg_stems) if svg_stems else 0
        check_pptx(pptx_path, expected, result)
        print()

    # Report
    print("-" * 50)
    result.report()
    print()

    total = len(result.passes) + len(result.warns) + len(result.errors)
    print(
        f"Result: {len(result.passes)} passed, {len(result.warns)} warnings, "
        f"{len(result.errors)} errors / {total} checks"
    )

    if result.has_errors:
        return 1
    if args.strict and result.warns:
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
